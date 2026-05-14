"""PPTx 리포트 생성 유틸 — Page 7의 종합 리포트를 PowerPoint로 변환."""
from __future__ import annotations

import io
from datetime import datetime
from typing import Iterable

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Inches, Pt

# 16:9 슬라이드 (Pt 기준)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 색상 팔레트
COLOR_PRIMARY    = RGBColor(0x4F, 0x46, 0xE5)
COLOR_SECONDARY  = RGBColor(0x7C, 0x3A, 0xED)
COLOR_DARK       = RGBColor(0x1E, 0x29, 0x3B)
COLOR_GRAY       = RGBColor(0x64, 0x74, 0x8B)
COLOR_LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
COLOR_BG         = RGBColor(0xF8, 0xFA, 0xFC)
COLOR_RED        = RGBColor(0xDC, 0x26, 0x26)
COLOR_AMBER      = RGBColor(0xD9, 0x77, 0x06)
COLOR_GREEN      = RGBColor(0x16, 0xA3, 0x4A)
COLOR_BLUE       = RGBColor(0x25, 0x63, 0xEB)


# ── 기본 헬퍼 ──────────────────────────────────────────────────────
def _set_slide_size(prs: Presentation):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def _add_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def _add_textbox(slide, left, top, width, height,
                  text: str, *, font_size: int = 14,
                  bold: bool = False, color: RGBColor = COLOR_DARK,
                  align: PP_ALIGN = PP_ALIGN.LEFT,
                  font_name: str = "Malgun Gothic"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top  = Cm(0.05); tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_filled_rect(slide, left, top, width, height,
                      fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_table(slide, left, top, width, height,
                df: pd.DataFrame, *, font_size: int = 10,
                header_color: RGBColor = COLOR_PRIMARY):
    """DataFrame을 PPT 테이블로 변환."""
    rows, cols = df.shape[0] + 1, df.shape[1]
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    # 헤더
    for ci, col_name in enumerate(df.columns):
        cell = table.cell(0, ci)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Malgun Gothic"
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 데이터
    for ri in range(df.shape[0]):
        for ci in range(df.shape[1]):
            cell = table.cell(ri + 1, ci)
            v = df.iloc[ri, ci]
            cell.text = "" if pd.isna(v) else str(v)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Malgun Gothic"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = COLOR_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                COLOR_BG if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            )

    return tbl_shape


def _add_header_bar(slide, title: str, subtitle: str = ""):
    """슬라이드 상단 헤더(보라 그라디언트 느낌의 단색 막대)."""
    _add_filled_rect(slide, 0, 0, SLIDE_W, Cm(2.0), COLOR_PRIMARY)
    _add_textbox(slide, Cm(1), Cm(0.4), SLIDE_W - Cm(2), Cm(0.9),
                  title, font_size=22, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        _add_textbox(slide, Cm(1), Cm(1.25), SLIDE_W - Cm(2), Cm(0.6),
                      subtitle, font_size=11,
                      color=RGBColor(0xE0, 0xE7, 0xFF))


def _add_footer(slide, page_num: int, total: int):
    _add_textbox(slide, SLIDE_W - Cm(4.5), SLIDE_H - Cm(0.7),
                  Cm(4), Cm(0.5),
                  f"{page_num} / {total}",
                  font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


# ── 슬라이드 빌더 ──────────────────────────────────────────────────
def _build_title_slide(prs: Presentation, kpi: dict):
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_PRIMARY)
    # 브랜드 라인
    _add_textbox(slide, Cm(2), Cm(1.5), SLIDE_W - Cm(4), Cm(0.8),
                  "LG DISPLAY", font_size=14, bold=True,
                  color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)
    # 메인 타이틀
    _add_textbox(slide, Cm(2), Cm(2.4), SLIDE_W - Cm(4), Cm(1.6),
                  "Customer Satisfaction Survey", font_size=36, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Cm(2), Cm(4.0), SLIDE_W - Cm(4), Cm(1.2),
                  "Analysis Platform — ML 종합 리포트", font_size=24, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Cm(2), Cm(5.3), SLIDE_W - Cm(4), Cm(0.8),
                  "EDA · Feature Engineering · Model Comparison · Feature Importance · 추가 과제 제안",
                  font_size=12, color=RGBColor(0xE0, 0xE7, 0xFF),
                  align=PP_ALIGN.CENTER)
    _add_textbox(slide, Cm(2), Cm(6.1), SLIDE_W - Cm(4), Cm(0.6),
                  f"생성일: {datetime.now():%Y-%m-%d %H:%M}",
                  font_size=11, color=RGBColor(0xC7, 0xD2, 0xFE),
                  align=PP_ALIGN.CENTER)
    # 하단 KPI
    box_w = (SLIDE_W - Cm(6)) / 4
    top = SLIDE_H - Cm(4.0)
    items = [
        ("최고 Macro F1", kpi.get("best_f1_val", 0), kpi.get("best_f1_model", "-")),
        ("최고 Macro AUROC", kpi.get("best_au_val", 0), kpi.get("best_au_model", "-")),
        ("학습 모델 수", kpi.get("n_models", 0), "models"),
        ("반복 시드", kpi.get("n_seeds", 0), "seeds"),
    ]
    for i, (label, val, sub) in enumerate(items):
        x = Cm(2) + i * (box_w + Cm(0.5))
        _add_filled_rect(slide, x, top, box_w, Cm(2.6),
                          RGBColor(0xFF, 0xFF, 0xFF))
        val_str = f"{val:.4f}" if isinstance(val, float) else str(val)
        _add_textbox(slide, x, top + Cm(0.2), box_w, Cm(0.6),
                      label, font_size=11, color=COLOR_GRAY,
                      align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, top + Cm(0.8), box_w, Cm(1.0),
                      val_str, font_size=24, bold=True, color=COLOR_PRIMARY,
                      align=PP_ALIGN.CENTER)
        _add_textbox(slide, x, top + Cm(1.9), box_w, Cm(0.5),
                      str(sub), font_size=10, color=COLOR_DARK,
                      align=PP_ALIGN.CENTER)


def _build_kpi_settings_slide(prs: Presentation, settings_df: pd.DataFrame, perf_df: pd.DataFrame):
    slide = _add_blank_slide(prs)
    _add_header_bar(slide, "분석 설정 & 모델 성능 종합",
                     "Configuration · Performance Summary")
    # 좌: 설정 테이블
    _add_textbox(slide, Cm(1), Cm(2.4), Cm(5), Cm(0.6),
                  "📌 분석 설정", font_size=14, bold=True, color=COLOR_DARK)
    _add_table(slide, Cm(1), Cm(3.0), Cm(11), Cm(4),
                settings_df, font_size=10)

    # 우: 성능 테이블
    _add_textbox(slide, Cm(13.5), Cm(2.4), Cm(10), Cm(0.6),
                  "📋 모델 성능 (Mean ± Std)", font_size=14, bold=True, color=COLOR_DARK)
    _add_table(slide, Cm(13.5), Cm(3.0), SLIDE_W - Cm(14.5), Cm(4),
                perf_df, font_size=10)


def _build_table_slide(prs: Presentation, title: str, subtitle: str,
                        df: pd.DataFrame, max_rows: int = 12,
                        accent: RGBColor = COLOR_PRIMARY):
    """제목 + 부가 설명 + 표만 들어가는 슬라이드."""
    slide = _add_blank_slide(prs)
    _add_header_bar(slide, title, subtitle)
    if df is None or df.empty:
        _add_textbox(slide, Cm(1), Cm(3), SLIDE_W - Cm(2), Cm(2),
                      "(분석 결과가 없습니다)",
                      font_size=14, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
        return
    show_df = df.head(max_rows).copy()
    # 너무 길면 텍스트 줄임
    for c in show_df.columns:
        show_df[c] = show_df[c].apply(
            lambda v: (str(v)[:60] + "...") if isinstance(v, str) and len(str(v)) > 60 else v
        )
    height = min(Cm(1.0) + Cm(0.45) * (len(show_df) + 1), Cm(13.5))
    _add_table(slide, Cm(1), Cm(2.4), SLIDE_W - Cm(2), height,
                show_df, font_size=10, header_color=accent)


def _build_insights_slide(prs: Presentation, title: str, insights: list[str],
                           accent: RGBColor = COLOR_PRIMARY):
    slide = _add_blank_slide(prs)
    _add_header_bar(slide, title, "Key Insights")
    box_top = Cm(2.4)
    box_h = SLIDE_H - Cm(3.5)
    _add_filled_rect(slide, Cm(1), box_top, SLIDE_W - Cm(2), box_h, COLOR_BG, COLOR_LIGHT_GRAY)
    tb = slide.shapes.add_textbox(Cm(1.5), box_top + Cm(0.3),
                                     SLIDE_W - Cm(3), box_h - Cm(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(insights):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = "Malgun Gothic"
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_DARK
        p.space_after = Pt(6)


def _build_tasks_slides(prs: Presentation, category_key: str, category: dict):
    """추가 과제 카테고리당 1슬라이드 — H/M/L 섹션 모두 포함."""
    slide = _add_blank_slide(prs)
    _add_header_bar(slide, f"{category['icon']} {category['title']}",
                     "추가 과제 제안 — Priority별 후속 작업")

    rows = []
    for level, label in [("high", "🔴 High"), ("mid", "🟡 Mid"), ("low", "🟢 Low")]:
        for it in category.get(level, []):
            rows.append({
                "우선순위": label,
                "과제": it["name"],
                "기대 효과": it["outcome"],
            })
    df_tasks = pd.DataFrame(rows)
    if df_tasks.empty:
        _add_textbox(slide, Cm(1), Cm(3), SLIDE_W - Cm(2), Cm(2),
                      "(추가 과제 정의 없음)", font_size=14,
                      color=COLOR_GRAY, align=PP_ALIGN.CENTER)
        return
    df_tasks["과제"] = df_tasks["과제"].apply(
        lambda v: v if len(str(v)) <= 50 else str(v)[:50] + "..."
    )
    df_tasks["기대 효과"] = df_tasks["기대 효과"].apply(
        lambda v: v if len(str(v)) <= 60 else str(v)[:60] + "..."
    )
    height = min(Cm(1.0) + Cm(0.5) * (len(df_tasks) + 1), Cm(14))
    _add_table(slide, Cm(0.7), Cm(2.4), SLIDE_W - Cm(1.4), height,
                df_tasks, font_size=10, header_color=COLOR_DARK)

    if "checklist" in category:
        # 별도 슬라이드 — 체크리스트
        cl_slide = _add_blank_slide(prs)
        _add_header_bar(cl_slide, "운영 모델 배포 체크리스트",
                         "Production Deployment Checklist")
        for i, item in enumerate(category["checklist"]):
            top = Cm(2.6) + i * Cm(0.7)
            _add_textbox(cl_slide, Cm(2), top, SLIDE_W - Cm(4), Cm(0.6),
                          f"☐  {item}", font_size=14, color=COLOR_DARK)


def _build_closing_slide(prs: Presentation, top_features_text: str):
    slide = _add_blank_slide(prs)
    _add_filled_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_DARK)
    _add_textbox(slide, Cm(2), Cm(2.5), SLIDE_W - Cm(4), Cm(1.5),
                  "Thank You", font_size=54, bold=True,
                  color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _add_textbox(slide, Cm(2), Cm(4.5), SLIDE_W - Cm(4), Cm(0.8),
                  "ML 분석 종합 리포트 — Streamlit Dashboard에서 생성",
                  font_size=14, color=RGBColor(0xC7, 0xD2, 0xFE),
                  align=PP_ALIGN.CENTER)
    if top_features_text and top_features_text != "—":
        _add_textbox(slide, Cm(2), Cm(5.6), SLIDE_W - Cm(4), Cm(0.6),
                      f"Top Features: {top_features_text}",
                      font_size=11, color=RGBColor(0xA5, 0xB4, 0xFC),
                      align=PP_ALIGN.CENTER)


# ── 메인 빌드 함수 ──────────────────────────────────────────────────
def build_pptx_report(*,
                       kpi: dict,
                       settings_df: pd.DataFrame,
                       perf_df: pd.DataFrame,
                       eda_kw_df: pd.DataFrame | None,
                       fe_info_df: pd.DataFrame,
                       friedman_df: pd.DataFrame,
                       wilcoxon_df: pd.DataFrame | None,
                       importance_df: pd.DataFrame | None,
                       shap_df: pd.DataFrame | None,
                       insights: dict,
                       additional_tasks: dict,
                       top_features_text: str = "—") -> bytes:
    """모든 섹션을 슬라이드로 렌더링 후 bytes 반환."""
    prs = Presentation()
    _set_slide_size(prs)

    # 1) 타이틀
    _build_title_slide(prs, kpi)

    # 2) 설정 + 성능 종합
    _build_kpi_settings_slide(prs, settings_df, perf_df)

    # 3) EDA 인사이트
    _build_insights_slide(prs, "🔍 EDA 결과 — 핵심 인사이트",
                            insights["eda"]["items"], COLOR_BLUE)

    # 4) EDA — KW 검정 결과
    if eda_kw_df is not None and not eda_kw_df.empty:
        kw_show = eda_kw_df.head(15)
        _build_table_slide(prs, "🔍 Kruskal-Wallis 단변량 검정 (Top 15)",
                            "각 변수의 클래스 분리 능력 (η² 큰 순)",
                            kw_show, max_rows=15, accent=COLOR_BLUE)

    # 5) Feature Engineering
    _build_table_slide(prs, "🛠 Feature Engineering 적용 전략",
                        "Strategy & Output Shape",
                        fe_info_df, max_rows=10, accent=COLOR_GREEN)
    _build_insights_slide(prs, "🛠 Feature Engineering — 핵심 인사이트",
                            insights["fe"]["items"], COLOR_GREEN)

    # 6) 모델 비교
    _build_table_slide(prs, "🤖 모델 비교 — Friedman Test",
                        "전체 모델 차이 유의성 검정",
                        friedman_df, max_rows=10, accent=COLOR_AMBER)
    if wilcoxon_df is not None and not wilcoxon_df.empty:
        _build_table_slide(prs, "🤖 Pairwise Wilcoxon Signed-Rank",
                            "Macro F1 기준 Bonferroni 보정",
                            wilcoxon_df, max_rows=12, accent=COLOR_AMBER)
    _build_insights_slide(prs, "🤖 모델 비교 — 핵심 인사이트",
                            insights["model"]["items"], COLOR_AMBER)

    # 7) Feature Importance
    if importance_df is not None and not importance_df.empty:
        _build_table_slide(prs, "🎯 Feature Importance Top-15",
                            "내장 또는 SHAP 기반",
                            importance_df.head(15), max_rows=15, accent=COLOR_RED)
    if shap_df is not None and not shap_df.empty:
        _build_table_slide(prs, "🎯 SHAP Top-15 (방향 포함)",
                            "양수=클래스↑ 기여 / 음수=감소 기여",
                            shap_df.head(15), max_rows=15, accent=COLOR_RED)
    _build_insights_slide(prs, "🎯 Feature Importance — 핵심 인사이트",
                            insights["fi"]["items"], COLOR_RED)

    # 8) 추가 과제 (4 카테고리 × 1슬라이드)
    for key in ["eda", "fe", "model", "fi"]:
        if key in additional_tasks:
            _build_tasks_slides(prs, key, additional_tasks[key])

    # 9) 종료
    _build_closing_slide(prs, top_features_text)

    # 페이지 번호 추가
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i == 0 or i == total - 1:
            continue
        _add_footer(slide, i + 1, total)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
