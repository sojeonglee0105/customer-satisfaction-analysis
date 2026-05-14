"""
CSAT Brief Slides 생성기
========================
3개 핵심 슬라이드(+ 표지/마무리)로 구성된 PPTX를 만든다.

  1) customer_satisfaction_full_columns.xlsx → Train/Test 분할 데이터셋 소개
     (무엇을 활용해서, 무엇을 예측하는지)
  2) 가장 성능이 좋은 세팅의 특징
  3) 그럼에도 불구하고 모델이 잘 못 맞히는 케이스

차트는 matplotlib으로 PNG 저장 후 PPTX에 임베드한다.
한글 폰트는 Windows 기본 'Malgun Gothic'을 사용한다.
"""
from __future__ import annotations

import warnings, sys
warnings.filterwarnings("ignore")
sys.stdout.reconfigure(encoding="utf-8")

from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams

from sklearn.preprocessing import LabelEncoder, StandardScaler
from sklearn.linear_model import LogisticRegression
from sklearn.tree import DecisionTreeClassifier
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import (
    f1_score, roc_auc_score, accuracy_score, confusion_matrix,
)
from sklearn.preprocessing import label_binarize
from lightgbm import LGBMClassifier
from xgboost import XGBClassifier

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt

# ── 경로 ────────────────────────────────────────────────────────────
BASE        = Path(__file__).resolve().parent.parent
DATA_DIR    = BASE / "data"
REPORT_DIR  = BASE / "reports"
CHART_DIR   = REPORT_DIR / "_brief_charts"
REPORT_DIR.mkdir(exist_ok=True, parents=True)
CHART_DIR.mkdir(exist_ok=True, parents=True)

TRAIN_CSV = DATA_DIR / "customer_satisfaction_train.csv"
TEST_CSV  = DATA_DIR / "customer_satisfaction_test.csv"
SOURCE_XLSX = "customer_satisfaction_full_columns.xlsx"

# ── LG Display CI 색상 팔레트 ───────────────────────────────────────
# LG Active Red(#A50034) Primary · Black/Grey Mono · White Bg
COLOR_PRIMARY    = RGBColor(0xA5, 0x00, 0x34)   # LG Active Red
COLOR_PRIMARY_DK = RGBColor(0x7A, 0x00, 0x26)   # Dark Red (gradient)
COLOR_PRIMARY_LT = RGBColor(0xE6, 0x00, 0x3C)   # Bright Red (highlight)
COLOR_SECONDARY  = RGBColor(0xC8, 0x10, 0x2E)   # Bright Red accent
COLOR_DARK       = RGBColor(0x1A, 0x1A, 0x1A)   # LG Black
COLOR_GRAY       = RGBColor(0x59, 0x59, 0x59)   # LG Mid Grey
COLOR_LIGHT_GRAY = RGBColor(0xD9, 0xD9, 0xD9)   # Border Grey
COLOR_BG         = RGBColor(0xF5, 0xF5, 0xF5)   # Soft Grey BG
COLOR_RED        = RGBColor(0xC8, 0x10, 0x2E)   # Alert Red
COLOR_AMBER      = RGBColor(0xD9, 0x77, 0x06)
COLOR_BLUE       = RGBColor(0x59, 0x59, 0x59)   # blue 자리는 회색으로
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_INDIGO_LT  = RGBColor(0xFA, 0xE5, 0xEB)   # Pale Pink (header subtitle)
COLOR_TITLE_BG   = RGBColor(0x1A, 0x1A, 0x1A)   # 표지 배경 — LG Black

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# matplotlib 한글
for fname in ["Malgun Gothic", "NanumGothic", "Apple SD Gothic Neo", "Arial"]:
    try:
        rcParams["font.family"] = fname
        break
    except Exception:
        pass
rcParams["axes.unicode_minus"] = False
rcParams["font.size"] = 11


# ===================================================================
# 1) 데이터 로드 & 분할 메타정보
# ===================================================================
def load_data():
    train = pd.read_csv(TRAIN_CSV)
    test  = pd.read_csv(TEST_CSV)
    return train, test


def dataset_meta(train: pd.DataFrame, test: pd.DataFrame) -> dict:
    full_n = len(train) + len(test)
    feat_cols = [c for c in train.columns if c not in ["year", "rpi"]]
    cat_cols = ["area", "product", "client"]
    num_cols = [c for c in feat_cols if c not in cat_cols]
    cls_train = train["rpi"].value_counts().sort_index()
    cls_test  = test["rpi"].value_counts().sort_index()
    return {
        "source_file": SOURCE_XLSX,
        "full_n": full_n,
        "n_train": len(train),
        "n_test":  len(test),
        "split_ratio": f"{len(train)/full_n:.0%} : {len(test)/full_n:.0%}",
        "n_features": len(feat_cols),
        "n_numeric":  len(num_cols),
        "n_categorical": len(cat_cols),
        "feat_cols": feat_cols,
        "cat_cols":  cat_cols,
        "num_cols":  num_cols,
        "years": sorted(set(train["year"]).union(set(test["year"]))),
        "areas": sorted(train["area"].unique()),
        "products": sorted(train["product"].unique()),
        "clients": sorted(train["client"].unique()),
        "cls_train": cls_train,
        "cls_test":  cls_test,
        "target": "rpi (Risk Performance Index, 1=Best ~ 5=Worst)",
    }


# ===================================================================
# 2) 모델 5종 학습 + Best 선정
# ===================================================================
def encode_and_split(train: pd.DataFrame, test: pd.DataFrame):
    cat_cols = ["area", "product", "client"]
    drop_cols = ["year", "rpi"]
    train, test = train.copy(), test.copy()
    for c in cat_cols:
        le = LabelEncoder()
        le.fit(pd.concat([train[c], test[c]]))
        train[c] = le.transform(train[c])
        test[c]  = le.transform(test[c])
    feat_cols = [c for c in train.columns if c not in drop_cols]
    X_train = train[feat_cols].values
    X_test  = test[feat_cols].values
    y_train = train["rpi"].values
    y_test  = test["rpi"].values
    sc = StandardScaler().fit(X_train)
    return (sc.transform(X_train), y_train, sc.transform(X_test), y_test, feat_cols)


def build_models():
    return {
        "Logistic Regression":
            LogisticRegression(max_iter=2000, class_weight="balanced", random_state=42, n_jobs=None),
        "Decision Tree":
            DecisionTreeClassifier(class_weight="balanced", random_state=42),
        "Random Forest":
            RandomForestClassifier(n_estimators=400, class_weight="balanced",
                                    random_state=42, n_jobs=-1),
        "LightGBM":
            LGBMClassifier(n_estimators=400, learning_rate=0.05,
                            class_weight="balanced", random_state=42,
                            verbosity=-1, n_jobs=-1),
        "XGBoost":
            XGBClassifier(n_estimators=400, learning_rate=0.05,
                           random_state=42, n_jobs=-1, verbosity=0,
                           use_label_encoder=False, eval_metric="mlogloss"),
    }


def evaluate(models: dict, X_tr, y_tr, X_te, y_te, all_classes):
    rows = []
    fitted = {}
    for name, m in models.items():
        if name == "XGBoost":
            sw = pd.Series(y_tr).map(
                {c: len(y_tr) / (len(np.unique(y_tr)) * (y_tr == c).sum())
                 for c in np.unique(y_tr)}
            ).values
            m.fit(X_tr, y_tr - 1, sample_weight=sw)
            yhat   = m.predict(X_te) + 1
            yproba = m.predict_proba(X_te)
        else:
            m.fit(X_tr, y_tr)
            yhat   = m.predict(X_te)
            yproba = m.predict_proba(X_te)
        # AUROC: full one-vs-rest using all training classes
        y_te_bin = label_binarize(y_te, classes=all_classes)
        try:
            macro_auroc = roc_auc_score(y_te_bin, yproba, average="macro", multi_class="ovr")
            micro_auroc = roc_auc_score(y_te_bin, yproba, average="micro", multi_class="ovr")
        except Exception:
            macro_auroc, micro_auroc = np.nan, np.nan
        rows.append({
            "Model": name,
            "Accuracy": accuracy_score(y_te, yhat),
            "Macro F1": f1_score(y_te, yhat, average="macro"),
            "Micro F1": f1_score(y_te, yhat, average="micro"),
            "Macro AUROC": macro_auroc,
            "Micro AUROC": micro_auroc,
        })
        fitted[name] = (m, yhat, yproba)
    return pd.DataFrame(rows), fitted


# ===================================================================
# 3) 오분류 분석
# ===================================================================
def misclass_analysis(test_df: pd.DataFrame, y_true, y_pred, y_proba):
    df = test_df.copy().reset_index(drop=True)
    df["pred"]   = y_pred
    df["true"]   = y_true
    df["correct"] = (df["pred"] == df["true"]).astype(int)
    df["pred_conf"] = y_proba.max(axis=1)
    df["pred_gap"]  = df["true"].astype(int) - df["pred"].astype(int)
    df["abs_gap"]   = df["pred_gap"].abs()

    # 그룹별 오류율
    err_by_class = (
        df.groupby("true")["correct"]
          .agg(["count", "sum"])
          .assign(error_rate=lambda d: 1 - d["sum"] / d["count"])
          .reset_index()
          .rename(columns={"true": "RPI(실제)", "count": "샘플수", "sum": "정답수"})
    )
    err_by_product = (
        df.groupby("product")["correct"]
          .agg(["count", "sum"])
          .assign(error_rate=lambda d: 1 - d["sum"] / d["count"])
          .reset_index()
          .rename(columns={"count": "샘플수", "sum": "정답수"})
          .sort_values("error_rate", ascending=False)
    )
    err_by_area = (
        df.groupby("area")["correct"]
          .agg(["count", "sum"])
          .assign(error_rate=lambda d: 1 - d["sum"] / d["count"])
          .reset_index()
          .rename(columns={"count": "샘플수", "sum": "정답수"})
          .sort_values("error_rate", ascending=False)
    )
    # 가장 어려운 케이스(예측 갭 큰 것 + 신뢰도 높은 오분류)
    hard = (
        df[df["correct"] == 0]
        .sort_values(["abs_gap", "pred_conf"], ascending=[False, False])
        .head(10)[["year", "area", "product", "client",
                    "true", "pred", "pred_conf", "abs_gap"]]
        .rename(columns={"true": "실제 RPI", "pred": "예측 RPI",
                          "pred_conf": "예측 신뢰도", "abs_gap": "오차 |Δ|"})
    )
    return df, err_by_class, err_by_product, err_by_area, hard


# ===================================================================
# 4) 차트 (matplotlib)
# ===================================================================
def chart_class_dist(meta: dict, out: Path):
    fig, ax = plt.subplots(figsize=(8, 4.2))
    classes = sorted(set(meta["cls_train"].index).union(meta["cls_test"].index))
    train_v = [int(meta["cls_train"].get(c, 0)) for c in classes]
    test_v  = [int(meta["cls_test"].get(c, 0))  for c in classes]
    x = np.arange(len(classes))
    w = 0.4
    # LG CI: Active Red (Train) + Mid Grey (Test)
    ax.bar(x - w/2, train_v, w, label=f"Train (n={meta['n_train']})", color="#A50034")
    ax.bar(x + w/2, test_v,  w, label=f"Test  (n={meta['n_test']})", color="#A0A0A0")
    for i, (a, b) in enumerate(zip(train_v, test_v)):
        ax.text(i - w/2, a + max(train_v) * 0.015, f"{a}", ha="center", fontsize=9)
        ax.text(i + w/2, b + max(train_v) * 0.015, f"{b}", ha="center", fontsize=9)
    ax.set_xticks(x); ax.set_xticklabels([f"RPI {c}" for c in classes])
    ax.set_ylabel("샘플 수")
    ax.set_title("RPI 클래스별 Train / Test 분포 (전체 1500행 → 70:30 분할)")
    ax.legend(loc="upper right")
    ax.grid(axis="y", alpha=0.25)
    fig.tight_layout(); fig.savefig(out, dpi=160); plt.close(fig)


def chart_perf_compare(perf_df: pd.DataFrame, out: Path):
    fig, ax = plt.subplots(figsize=(9, 4.5))
    metrics = ["Macro F1", "Micro F1", "Macro AUROC", "Micro AUROC"]
    # LG CI red gradation + greys
    colors  = ["#A50034", "#E6003C", "#595959", "#B0B0B0"]
    x = np.arange(len(perf_df))
    w = 0.18
    for i, (m, c) in enumerate(zip(metrics, colors)):
        ax.bar(x + (i - 1.5) * w, perf_df[m].values, w, label=m, color=c)
    ax.set_xticks(x); ax.set_xticklabels(perf_df["Model"], rotation=15)
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("Score")
    ax.set_title("모델별 성능 비교 (Class-Weighted, 평가용 1방 학습)")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(loc="lower right", ncols=4, fontsize=9)
    fig.tight_layout(); fig.savefig(out, dpi=160); plt.close(fig)


def chart_confusion(y_true, y_pred, classes, model_name: str, out: Path):
    cm = confusion_matrix(y_true, y_pred, labels=classes)
    cm_norm = cm.astype(float) / cm.sum(axis=1, keepdims=True).clip(min=1)
    fig, ax = plt.subplots(figsize=(5.5, 4.7))
    # LG Active Red gradient colormap
    from matplotlib.colors import LinearSegmentedColormap
    lg_cmap = LinearSegmentedColormap.from_list(
        "LG_Red", ["#FFFFFF", "#FAE5EB", "#E6003C", "#A50034", "#7A0026"], N=256)
    im = ax.imshow(cm_norm, cmap=lg_cmap, vmin=0, vmax=1)
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            v = cm_norm[i, j]
            ax.text(j, i, f"{cm[i,j]}\n({v:.0%})",
                    ha="center", va="center", fontsize=9,
                    color="white" if v > 0.55 else "#1E293B")
    ax.set_xticks(range(len(classes))); ax.set_yticks(range(len(classes)))
    ax.set_xticklabels([f"P{c}" for c in classes])
    ax.set_yticklabels([f"T{c}" for c in classes])
    ax.set_xlabel("예측 RPI"); ax.set_ylabel("실제 RPI")
    ax.set_title(f"Confusion Matrix — {model_name}")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    fig.tight_layout(); fig.savefig(out, dpi=160); plt.close(fig)


def chart_error_by_group(df: pd.DataFrame, group_col: str, label: str, out: Path):
    fig, ax = plt.subplots(figsize=(7.5, 4))
    show = df.head(8).copy()
    # LG Active Red(top) + Mid Grey(others)
    colors = ["#A50034" if v >= show["error_rate"].max() * 0.9 else "#B0B0B0"
              for v in show["error_rate"]]
    ax.barh(show[group_col].astype(str), show["error_rate"] * 100, color=colors)
    ax.invert_yaxis()
    for i, (rate, n) in enumerate(zip(show["error_rate"], show["샘플수"])):
        ax.text(rate * 100 + 0.4, i, f"{rate:.1%} (n={n})",
                va="center", fontsize=9, color="#1E293B")
    ax.set_xlabel("오류율 (%)")
    ax.set_title(f"{label}별 오류율 — Top 8")
    ax.set_xlim(0, max(50, show["error_rate"].max() * 100 * 1.3))
    ax.grid(axis="x", alpha=0.2)
    fig.tight_layout(); fig.savefig(out, dpi=160); plt.close(fig)


# ===================================================================
# 5) PPTX 빌더
# ===================================================================
def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def text(slide, left, top, w, h, txt, *, fs=14, bold=False,
         color=COLOR_DARK, align=PP_ALIGN.LEFT, font="Malgun Gothic"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
    tf.margin_top  = Cm(0.05); tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.name = font; r.font.size = Pt(fs)
    r.font.bold = bold; r.font.color.rgb = color
    return tb


def rect(slide, left, top, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def table_df(slide, left, top, w, h, df: pd.DataFrame, *,
              fs=10, header=COLOR_PRIMARY, first_col_bold=False):
    rows, cols = df.shape[0] + 1, df.shape[1]
    tbl = slide.shapes.add_table(rows, cols, left, top, w, h).table
    for ci, name in enumerate(df.columns):
        c = tbl.cell(0, ci); c.text = str(name)
        c.fill.solid(); c.fill.fore_color.rgb = header
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Malgun Gothic"
                r.font.size = Pt(fs); r.font.bold = True
                r.font.color.rgb = COLOR_WHITE
    for ri in range(df.shape[0]):
        for ci in range(df.shape[1]):
            cell = tbl.cell(ri + 1, ci)
            v = df.iloc[ri, ci]
            cell.text = "" if pd.isna(v) else str(v)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "Malgun Gothic"
                    r.font.size = Pt(fs)
                    r.font.color.rgb = COLOR_DARK
                    if first_col_bold and ci == 0:
                        r.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG if ri % 2 == 0 else COLOR_WHITE
    return tbl


def header_bar(slide, title, subtitle=""):
    """LG 스타일 헤더 — 검은 바 + 좌측 빨간 액센트 라인."""
    rect(slide, 0, 0, SLIDE_W, Cm(2.0), COLOR_DARK)
    rect(slide, 0, 0, Cm(0.4), Cm(2.0), COLOR_PRIMARY)
    text(slide, Cm(0.9), Cm(0.3), SLIDE_W - Cm(2), Cm(0.95),
         title, fs=22, bold=True, color=COLOR_WHITE)
    if subtitle:
        text(slide, Cm(0.9), Cm(1.25), SLIDE_W - Cm(2), Cm(0.6),
             subtitle, fs=11, color=COLOR_INDIGO_LT)
    # 우측 LG 브랜드 마크
    rect(slide, SLIDE_W - Cm(2.3), Cm(0.55), Cm(1.7), Cm(0.9),
         COLOR_PRIMARY)
    text(slide, SLIDE_W - Cm(2.3), Cm(0.55), Cm(1.7), Cm(0.9),
         "LG", fs=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


def kpi_card(slide, x, y, w, h, label, value, sub):
    """LG 스타일 KPI 카드 — 흰 배경 + 상단 빨간 라인."""
    rect(slide, x, y, w, h, COLOR_WHITE, COLOR_LIGHT_GRAY)
    rect(slide, x, y, w, Cm(0.15), COLOR_PRIMARY)  # 상단 빨간 라인
    text(slide, x, y + Cm(0.3), w, Cm(0.6),
         label, fs=11, color=COLOR_GRAY, align=PP_ALIGN.CENTER)
    text(slide, x, y + Cm(0.9), w, Cm(1.0),
         str(value), fs=22, bold=True, color=COLOR_PRIMARY,
         align=PP_ALIGN.CENTER)
    text(slide, x, y + Cm(1.95), w, Cm(0.5),
         str(sub), fs=10, color=COLOR_DARK, align=PP_ALIGN.CENTER)


# ── 슬라이드 1: Title (LG CI 시그니처 디자인) ─────────────
def build_title(prs, kpi: dict):
    s = add_blank(prs)
    # 1) 배경: LG Black
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_TITLE_BG)
    # 2) 우상단 LG 시그니처 — 빨간 원 + LG
    rect(s, SLIDE_W - Cm(3.5), Cm(0.8), Cm(2.5), Cm(2.5),
         COLOR_PRIMARY)
    text(s, SLIDE_W - Cm(3.5), Cm(0.95), Cm(2.5), Cm(2.2),
         "LG", fs=44, bold=True, color=COLOR_WHITE,
         align=PP_ALIGN.CENTER)
    # 3) 좌측 빨간 세로 액센트
    rect(s, 0, 0, Cm(0.6), SLIDE_H, COLOR_PRIMARY)
    # 4) 브랜드 라인
    text(s, Cm(2), Cm(1.4), SLIDE_W - Cm(6), Cm(0.8),
         "LG DISPLAY", fs=14, bold=True,
         color=COLOR_PRIMARY_LT, align=PP_ALIGN.LEFT)
    text(s, Cm(2), Cm(2.0), SLIDE_W - Cm(6), Cm(0.6),
         "Innovation for a Better Life", fs=11,
         color=RGBColor(0xB0, 0xB0, 0xB0), align=PP_ALIGN.LEFT)

    # 5) 메인 타이틀
    text(s, Cm(2), Cm(3.2), SLIDE_W - Cm(4), Cm(1.6),
         "Customer Satisfaction Survey", fs=40, bold=True,
         color=COLOR_WHITE, align=PP_ALIGN.LEFT)
    text(s, Cm(2), Cm(4.5), SLIDE_W - Cm(4), Cm(1.2),
         "Brief — 데이터셋 · Best 모델 · 잘 못 맞히는 케이스", fs=20, bold=True,
         color=COLOR_PRIMARY_LT, align=PP_ALIGN.LEFT)

    # 6) 빨간 구분 라인
    rect(s, Cm(2), Cm(5.7), Cm(8), Cm(0.08), COLOR_PRIMARY)

    # 7) 생성일
    text(s, Cm(2), Cm(5.9), SLIDE_W - Cm(4), Cm(0.7),
         f"생성일: {datetime.now():%Y-%m-%d %H:%M}",
         fs=11, color=RGBColor(0xC0, 0xC0, 0xC0), align=PP_ALIGN.LEFT)

    # 8) 하단 KPI 4종 — 어두운 배경에 어울리게
    box_w = (SLIDE_W - Cm(6)) / 4
    top = SLIDE_H - Cm(3.4)
    items = [
        ("최고 Macro F1", f"{kpi['best_f1_val']:.3f}", kpi["best_f1_model"]),
        ("최고 Macro AUROC", f"{kpi['best_au_val']:.3f}", kpi["best_au_model"]),
        ("Train / Test", f"{kpi['n_train']} : {kpi['n_test']}", "70 : 30 split"),
        ("입력 변수", f"{kpi['n_features']}",
         f"수치 {kpi['n_numeric']} + 범주 {kpi['n_categorical']}"),
    ]
    for i, (lab, val, sub) in enumerate(items):
        x = Cm(2) + i * (box_w + Cm(0.5))
        # 어두운 배경 KPI 카드
        rect(s, x, top, box_w, Cm(2.6), RGBColor(0x2A, 0x2A, 0x2A),
             RGBColor(0x40, 0x40, 0x40))
        rect(s, x, top, box_w, Cm(0.15), COLOR_PRIMARY)
        text(s, x, top + Cm(0.3), box_w, Cm(0.6),
             lab, fs=11, color=RGBColor(0xC0, 0xC0, 0xC0),
             align=PP_ALIGN.CENTER)
        text(s, x, top + Cm(0.9), box_w, Cm(1.0),
             str(val), fs=22, bold=True, color=COLOR_PRIMARY_LT,
             align=PP_ALIGN.CENTER)
        text(s, x, top + Cm(1.95), box_w, Cm(0.5),
             str(sub), fs=10, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # 9) 푸터 슬로건
    text(s, Cm(2), SLIDE_H - Cm(0.8), SLIDE_W - Cm(4), Cm(0.5),
         "Life's Good   ·   LGD Customer Satisfaction Analysis Platform",
         fs=10, color=COLOR_PRIMARY_LT, align=PP_ALIGN.CENTER)


# ── 슬라이드 2: Dataset 소개 ──────────────────────────────
def build_dataset(prs, meta: dict, chart_path: Path):
    s = add_blank(prs)
    header_bar(s, "1. 데이터셋 소개 — 무엇을 활용해서, 무엇을 예측하는가",
               f"Source: {meta['source_file']}  →  Train/Test 분할 ({meta['split_ratio']})")

    # 좌측: 데이터셋 메타 + 입력/타깃 정의
    left_x = Cm(0.8); left_w = Cm(11.5)
    text(s, left_x, Cm(2.4), left_w, Cm(0.6),
         "📌 데이터 소스 & 분할", fs=14, bold=True, color=COLOR_DARK)
    meta_df = pd.DataFrame({
        "항목": ["원본 파일", "전체 샘플", "Train", "Test",
                "분할 비율", "수집 연도", "결측치"],
        "내용": [
            meta["source_file"],
            f"{meta['full_n']:,} 행",
            f"{meta['n_train']:,} 행",
            f"{meta['n_test']:,} 행",
            meta["split_ratio"],
            ", ".join(str(y) for y in meta["years"]),
            "0 (없음)",
        ],
    })
    table_df(s, left_x, Cm(3.0), left_w, Cm(4.4),
              meta_df, fs=11, first_col_bold=True)

    text(s, left_x, Cm(7.7), left_w, Cm(0.6),
         "🎯 입력 변수 & 예측 대상", fs=14, bold=True, color=COLOR_DARK)
    feat_df = pd.DataFrame({
        "구분": ["입력 X (53→52)", "수치형 (48)", "범주형 (3)", "예측 Y"],
        "설명": [
            "고객 만족도/불만족 응답 + 메타(연도/제품/고객/평가영역)",
            "도메인×지표×응답형: t1/t2/c/d/q1/q2 × CSI/CCI × res/core/comm/total",
            "area(평가영역 6) · product(제품 5) · client(고객 5)",
            "rpi (1=Best ~ 5=Worst, 5-class 다중분류)",
        ],
    })
    table_df(s, left_x, Cm(8.3), left_w, Cm(5.4), feat_df, fs=10)

    # 우측: 분포 차트
    right_x = left_x + left_w + Cm(0.6); right_w = SLIDE_W - right_x - Cm(0.8)
    text(s, right_x, Cm(2.4), right_w, Cm(0.6),
         "📊 RPI 클래스 분포 (Train vs Test)",
         fs=14, bold=True, color=COLOR_DARK)
    s.shapes.add_picture(str(chart_path), right_x, Cm(3.0),
                          width=right_w, height=Cm(7.6))

    # 도메인 키 카드
    text(s, right_x, Cm(11.0), right_w, Cm(0.6),
         "🔖 평가 도메인 (area) 키", fs=12, bold=True, color=COLOR_DARK)
    rect(s, right_x, Cm(11.6), right_w, Cm(2.0), COLOR_BG, COLOR_LIGHT_GRAY)
    text(s, right_x + Cm(0.3), Cm(11.7), right_w - Cm(0.6), Cm(2.0),
         "t1 신기술  ·  t2 개발  ·  c Cost  ·  d 공급  ·  "
         "q1 품질  ·  q2 서비스",
         fs=11, color=COLOR_DARK)


# ── 슬라이드 3: Best 세팅 특징 ──────────────────────────
def build_best_setting(prs, perf_df: pd.DataFrame, best_name: str,
                        feat_imp: pd.DataFrame, chart_path: Path,
                        cm_path: Path):
    s = add_blank(prs)
    header_bar(s, "2. 가장 성능이 좋은 세팅의 특징",
               "Best Configuration · Class-Weighted · 5-class 다중분류")

    # KPI 헤드라인
    best_row = perf_df[perf_df["Model"] == best_name].iloc[0]
    box_w = Cm(5.5); top = Cm(2.4); gap = Cm(0.4)
    kpis = [
        ("Best Model", best_name, "class_weight='balanced'"),
        ("Macro F1", f"{best_row['Macro F1']:.3f}",
         f"Micro F1 {best_row['Micro F1']:.3f}"),
        ("Macro AUROC", f"{best_row['Macro AUROC']:.3f}",
         f"Micro AUROC {best_row['Micro AUROC']:.3f}"),
        ("Accuracy", f"{best_row['Accuracy']:.3f}", "1-fold holdout"),
    ]
    for i, (lab, v, sub) in enumerate(kpis):
        x = Cm(0.8) + i * (box_w + gap)
        kpi_card(s, x, top, box_w, Cm(2.6), lab, v, sub)

    # 좌하: 모델별 비교 차트
    text(s, Cm(0.8), Cm(5.4), Cm(13), Cm(0.6),
         "📈 모델별 성능 비교", fs=14, bold=True, color=COLOR_DARK)
    s.shapes.add_picture(str(chart_path), Cm(0.8), Cm(6.0),
                          width=Cm(15), height=Cm(7.4))

    # 우상: 핵심 특징 박스
    fx = Cm(16.2); fw = SLIDE_W - fx - Cm(0.8)
    text(s, fx, Cm(5.4), fw, Cm(0.6),
         "🔑 Best 세팅의 핵심 특징", fs=14, bold=True, color=COLOR_DARK)
    rect(s, fx, Cm(6.0), fw, Cm(7.4), COLOR_BG, COLOR_LIGHT_GRAY)
    feat_lines = [
        f"• 모델: {best_name} — 비선형/상호작용 학습에 강함",
        "• 표준화(StandardScaler) 후 학습 — 스케일 다른 CSI/CCI 동시 처리",
        "• 모든 모델 class_weight='balanced'로 클래스 불균형 자동 보정 (RPI 5는 7%만)",
        "• Macro/Micro 지표 동시 추적 — 소수 클래스(RPI 4·5) 누락 방지",
        f"• 입력 53개 그대로 사용(범주형 LabelEncoded, year 제외) → 해석 가능성 ↑",
        "• Top 영향 변수에 q2(서비스)·t1(신기술)·c(Cost) 도메인 핵심 지표 포함",
    ]
    tb = s.shapes.add_textbox(fx + Cm(0.3), Cm(6.2), fw - Cm(0.6), Cm(7.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(feat_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = "Malgun Gothic"; r.font.size = Pt(11)
        r.font.color.rgb = COLOR_DARK
        p.space_after = Pt(6)


# ── 슬라이드 4: 잘 못 맞히는 케이스 ──────────────────────
def build_misclass(prs, hard_df: pd.DataFrame, err_class_df: pd.DataFrame,
                    err_product_df: pd.DataFrame, err_area_df: pd.DataFrame,
                    cm_path: Path, prod_path: Path, area_path: Path,
                    overall_acc: float):
    s = add_blank(prs)
    header_bar(s, "3. 그럼에도 모델이 잘 못 맞히는 케이스",
               "Confusion Matrix · 그룹별 오류율 · Hard Cases")

    # 상단 KPI
    box_w = Cm(5.5); top = Cm(2.4); gap = Cm(0.4)
    rpi_top1 = err_class_df.sort_values("error_rate", ascending=False).iloc[0]
    prod_top1 = err_product_df.iloc[0]
    area_top1 = err_area_df.iloc[0]
    kpis = [
        ("전체 정확도", f"{overall_acc:.1%}", "Test 450건"),
        ("가장 어려운 RPI", f"RPI {int(rpi_top1['RPI(실제)'])}",
         f"오류 {rpi_top1['error_rate']:.0%}  (n={int(rpi_top1['샘플수'])})"),
        ("가장 어려운 제품", str(prod_top1["product"]),
         f"오류 {prod_top1['error_rate']:.0%}  (n={int(prod_top1['샘플수'])})"),
        ("가장 어려운 도메인", str(area_top1["area"]),
         f"오류 {area_top1['error_rate']:.0%}  (n={int(area_top1['샘플수'])})"),
    ]
    for i, (lab, v, sub) in enumerate(kpis):
        x = Cm(0.8) + i * (box_w + gap)
        kpi_card(s, x, top, box_w, Cm(2.6), lab, v, sub)

    # 좌: Confusion Matrix
    text(s, Cm(0.8), Cm(5.4), Cm(11.5), Cm(0.6),
         "🟦 Confusion Matrix", fs=14, bold=True, color=COLOR_DARK)
    s.shapes.add_picture(str(cm_path), Cm(0.8), Cm(6.0),
                          width=Cm(11.5), height=Cm(8.5))

    # 우상: Top 어려운 그룹 (제품, 도메인)
    rx = Cm(12.6); rw = SLIDE_W - rx - Cm(0.8)
    text(s, rx, Cm(5.4), rw, Cm(0.6),
         "🔥 그룹별 오류율 — 제품 / 평가 도메인",
         fs=14, bold=True, color=COLOR_DARK)
    s.shapes.add_picture(str(prod_path), rx, Cm(6.0),
                          width=rw, height=Cm(4.0))
    s.shapes.add_picture(str(area_path), rx, Cm(10.2),
                          width=rw, height=Cm(4.0))

    # 우하: 인사이트 박스
    text(s, rx, Cm(14.4), rw, Cm(0.6),
         "💡 시사점", fs=14, bold=True, color=COLOR_DARK)
    rect(s, rx, Cm(15.0), rw, Cm(3.2), COLOR_BG, COLOR_LIGHT_GRAY)
    insights = [
        "• 중간 등급(RPI 3·4)에서 인접 클래스로 1단계 오류가 다수 — 5단계 척도 본질의 모호성",
        "• 소수 클래스(RPI 5)는 표본 부족(78개)이라 confidence는 높아도 실제는 다른 등급으로 갈 수 있음",
        "• 특정 제품·도메인 조합(예: 자동차 d공급, 모니터 c비용)에서 오류 집중 → 도메인 맞춤 모델 또는 추가 피처 필요",
    ]
    tb = s.shapes.add_textbox(rx + Cm(0.3), Cm(15.1), rw - Cm(0.6), Cm(3.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(insights):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = "Malgun Gothic"; r.font.size = Pt(10.5)
        r.font.color.rgb = COLOR_DARK
        p.space_after = Pt(4)

    # 다음 슬라이드: Hard Cases Top 10 표
    s2 = add_blank(prs)
    header_bar(s2, "3-부록. 신뢰도 높지만 틀린 Hard Cases — Top 10",
               "|Δ| = |실제 RPI - 예측 RPI|, 신뢰도 = max(p)")
    show = hard_df.copy()
    show["예측 신뢰도"] = show["예측 신뢰도"].apply(lambda v: f"{v:.2f}")
    table_df(s2, Cm(0.8), Cm(2.4), SLIDE_W - Cm(1.6), Cm(11),
              show, fs=11, header=COLOR_RED)
    text(s2, Cm(0.8), Cm(13.6), SLIDE_W - Cm(1.6), Cm(1.4),
         "→ 예측 신뢰도가 높음에도 |Δ|가 2 이상인 케이스는 모델이 '확신을 가지고 틀리는' "
         "구조적 한계 영역. 도메인 전문가와의 정성적 검토 또는 추가 변수(B2B 거래조건, "
         "전년동기 대비 변동성 등) 도입을 권장합니다.",
         fs=11, color=COLOR_GRAY)


# ── 슬라이드 5: Closing (LG CI 스타일) ───────────────────
def build_closing(prs, best_name: str, top_features: list[str]):
    s = add_blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_TITLE_BG)
    # 좌측 빨간 액센트
    rect(s, 0, 0, Cm(0.6), SLIDE_H, COLOR_PRIMARY)
    # 우상단 LG 시그니처
    rect(s, SLIDE_W - Cm(3.5), Cm(0.8), Cm(2.5), Cm(2.5), COLOR_PRIMARY)
    text(s, SLIDE_W - Cm(3.5), Cm(0.95), Cm(2.5), Cm(2.2),
         "LG", fs=44, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    text(s, Cm(2), Cm(2.6), SLIDE_W - Cm(4), Cm(1.6),
         "Thank You", fs=54, bold=True,
         color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    rect(s, SLIDE_W / 2 - Cm(4), Cm(4.4), Cm(8), Cm(0.08), COLOR_PRIMARY)
    text(s, Cm(2), Cm(4.7), SLIDE_W - Cm(4), Cm(0.8),
         f"Best Model: {best_name}",
         fs=18, color=COLOR_PRIMARY_LT, align=PP_ALIGN.CENTER)
    if top_features:
        text(s, Cm(2), Cm(5.6), SLIDE_W - Cm(4), Cm(0.8),
             "Top Features: " + ", ".join(top_features[:6]),
             fs=12, color=RGBColor(0xC0, 0xC0, 0xC0),
             align=PP_ALIGN.CENTER)
    text(s, Cm(2), Cm(6.5), SLIDE_W - Cm(4), Cm(0.6),
         "LGD's Customer Satisfaction Survey Analysis Platform",
         fs=11, color=RGBColor(0xA0, 0xA0, 0xA0), align=PP_ALIGN.CENTER)
    text(s, Cm(2), SLIDE_H - Cm(0.8), SLIDE_W - Cm(4), Cm(0.5),
         "Life's Good", fs=11, bold=True, color=COLOR_PRIMARY_LT,
         align=PP_ALIGN.CENTER)


# ── 슬라이드 6: Reflection (느낀점 — 빈 슬라이드) ────────
def build_reflection(prs):
    s = add_blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_WHITE)
    # 헤더
    rect(s, 0, 0, SLIDE_W, Cm(2.0), COLOR_DARK)
    rect(s, 0, 0, Cm(0.4), Cm(2.0), COLOR_PRIMARY)
    text(s, Cm(0.9), Cm(0.3), SLIDE_W - Cm(2), Cm(0.95),
         "느낀점 (Reflection)", fs=22, bold=True, color=COLOR_WHITE)
    text(s, Cm(0.9), Cm(1.25), SLIDE_W - Cm(2), Cm(0.6),
         "본 분석을 진행하며 배운 점 / 인상 깊었던 점을 한 문장으로",
         fs=11, color=COLOR_INDIGO_LT)
    rect(s, SLIDE_W - Cm(2.3), Cm(0.55), Cm(1.7), Cm(0.9), COLOR_PRIMARY)
    text(s, SLIDE_W - Cm(2.3), Cm(0.55), Cm(1.7), Cm(0.9),
         "LG", fs=18, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    # 본문 — 큰 인용부호 카드 (빈 입력 영역)
    box_x = Cm(3); box_y = Cm(3.5)
    box_w = SLIDE_W - Cm(6); box_h = Cm(8.5)
    rect(s, box_x, box_y, box_w, box_h, COLOR_BG, COLOR_LIGHT_GRAY)
    # 좌상단 빨간 액센트 라인
    rect(s, box_x, box_y, box_w, Cm(0.15), COLOR_PRIMARY)
    # 큰 인용부호 (장식)
    text(s, box_x + Cm(0.4), box_y + Cm(0.2), Cm(2), Cm(2),
         "“", fs=72, bold=True, color=COLOR_PRIMARY_LT,
         align=PP_ALIGN.LEFT)
    text(s, box_x + box_w - Cm(2.2), box_y + box_h - Cm(2.2), Cm(2), Cm(2),
         "”", fs=72, bold=True, color=COLOR_PRIMARY_LT,
         align=PP_ALIGN.RIGHT)
    # 가이드 placeholder (회색)
    text(s, box_x + Cm(2), box_y + box_h / 2 - Cm(0.6),
         box_w - Cm(4), Cm(1.2),
         "여기에 한 문장을 직접 입력해 주세요",
         fs=14, color=RGBColor(0xB0, 0xB0, 0xB0),
         align=PP_ALIGN.CENTER)

    # 푸터
    rect(s, 0, SLIDE_H - Cm(1.0), SLIDE_W, Cm(1.0),
         RGBColor(0xF5, 0xF5, 0xF5))
    rect(s, 0, SLIDE_H - Cm(1.0), Cm(0.4), Cm(1.0), COLOR_PRIMARY)
    text(s, Cm(0.9), SLIDE_H - Cm(0.85), SLIDE_W - Cm(2), Cm(0.7),
         "LG DISPLAY  ·  Customer Satisfaction Survey Analysis Platform",
         fs=10, color=COLOR_GRAY, align=PP_ALIGN.LEFT)


# ===================================================================
# MAIN
# ===================================================================
def main():
    print("[1/6] 데이터 로드…")
    train, test = load_data()
    meta = dataset_meta(train, test)
    print(f"      train={meta['n_train']}  test={meta['n_test']}  feat={meta['n_features']}")

    print("[2/6] 모델 학습 (5종, class_weight='balanced')…")
    X_tr, y_tr, X_te, y_te, feat_cols = encode_and_split(train, test)
    all_classes = sorted(np.unique(y_tr))
    perf_df, fitted = evaluate(build_models(), X_tr, y_tr, X_te, y_te, all_classes)
    perf_df = perf_df.sort_values("Macro F1", ascending=False).reset_index(drop=True)
    print(perf_df.to_string(index=False))

    # Best 선정
    best_f1_row = perf_df.iloc[0]
    best_au_row = perf_df.sort_values("Macro AUROC", ascending=False).iloc[0]
    best_name = best_f1_row["Model"]
    best_model, best_yhat, best_proba = fitted[best_name]
    print(f"      → Best: {best_name}")

    print("[3/6] 오분류 분석…")
    df_full, err_class, err_prod, err_area, hard_df = misclass_analysis(
        test, y_te, best_yhat, best_proba)

    # Feature importance (가능한 경우)
    feat_imp = pd.DataFrame()
    if hasattr(best_model, "feature_importances_"):
        imps = best_model.feature_importances_
        feat_imp = pd.DataFrame({"feature": feat_cols, "importance": imps}) \
            .sort_values("importance", ascending=False).head(10)
    elif hasattr(best_model, "coef_"):
        imps = np.abs(best_model.coef_).mean(axis=0)
        feat_imp = pd.DataFrame({"feature": feat_cols, "importance": imps}) \
            .sort_values("importance", ascending=False).head(10)

    print("[4/6] 차트 생성…")
    cls_chart  = CHART_DIR / "class_dist.png"
    perf_chart = CHART_DIR / "perf_compare.png"
    cm_chart   = CHART_DIR / "confusion_matrix.png"
    prod_chart = CHART_DIR / "err_product.png"
    area_chart = CHART_DIR / "err_area.png"
    chart_class_dist(meta, cls_chart)
    chart_perf_compare(perf_df, perf_chart)
    chart_confusion(y_te, best_yhat, all_classes, best_name, cm_chart)
    chart_error_by_group(err_prod, "product", "제품", prod_chart)
    chart_error_by_group(err_area, "area", "평가 도메인", area_chart)

    print("[5/6] PPTX 빌드…")
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H

    kpi = {
        "best_f1_val":   best_f1_row["Macro F1"],
        "best_f1_model": best_name,
        "best_au_val":   best_au_row["Macro AUROC"],
        "best_au_model": best_au_row["Model"],
        "n_train": meta["n_train"], "n_test": meta["n_test"],
        "n_features": meta["n_features"],
        "n_numeric":  meta["n_numeric"],
        "n_categorical": meta["n_categorical"],
    }

    build_title(prs, kpi)
    build_dataset(prs, meta, cls_chart)
    build_best_setting(prs, perf_df.round(3), best_name, feat_imp,
                        perf_chart, cm_chart)
    build_misclass(prs, hard_df, err_class, err_prod, err_area,
                    cm_chart, prod_chart, area_chart,
                    overall_acc=best_f1_row["Accuracy"])
    build_closing(prs, best_name,
                   top_features=feat_imp["feature"].tolist() if not feat_imp.empty else [])
    build_reflection(prs)

    print("[6/6] 저장…")
    out = REPORT_DIR / f"CSAT_Brief_Slides_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    prs.save(out)
    print(f"\n✅ 저장 완료: {out}")
    print(f"   슬라이드 수: {len(prs.slides)}")
    return out


if __name__ == "__main__":
    main()
